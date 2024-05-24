import regex as re
import sys
from pptx import Presentation


books = {
  "kejadian": 1,
  "keluaran": 2,
  "imamat": 3,
  "bilangan": 4,
  "ulangan": 5,
  "yosua": 6,
  "hakim-hakim": 7,
  "rut": 8,
  "1 samuel": 9,
  "2 samuel": 10,
  "1 raja-raja": 11,
  "2 raja-raja": 12,
  "1 tawarikh": 13,
  "2 tawarikh": 14,
  "ezra": 15,
  "nehemia": 16,
  "ester": 17,
  "ayub": 18,
  "mazmur": 19,
  "amsal": 20,
  "pengkhotbah": 21,
  "kidung agung": 22,
  "yesaya": 23,
  "yeremia": 24,
  "ratapan": 25,
  "yehezkiel": 26,
  "daniel": 27,
  "hosea": 28,
  "yoel": 29,
  "amos": 30,
  "obaja": 31,
  "yunus": 32,
  "mikha": 33,
  "nahum": 34,
  "habakuk": 35,
  "zefanya": 36,
  "hagai": 37,
  "zakharia": 38,
  "maleakhi": 39,
  "matius": 40,
  "markus": 41,
  "lukas": 42,
  "yohanes": 43,
  "rasul": 44,
  "roma": 45,
  "i korintus": 46,
  "ii korintus": 47,
  "galatia": 48,
  "efesus": 49,
  "filipi": 50,
  "kolose": 51,
  "i tesalonika": 52,
  "ii tesalonika": 53,
  "i timotius": 54,
  "ii timotius": 55,
  "titus": 56,
  "filemon": 57,
  "ibrani": 58,
  "yakobus": 59,
  "i petrus": 60,
  "ii petrus": 61,
  "i yohanes": 62,
  "ii yohanes": 63,
  "iii yohanes": 64,
  "yudas": 65,
  "wahyu": 66
}


def extractBible(text):
    allBiblePattern = r"((?<=\\n|\s|\')(i{0,2}\s)?[a-z]+)\s*([0-9]+\s*\:\s*([0-9]+(\-|\,|(a|b|\:|\s))*)*)"
    bookPattern = r"^(?!kj|pkj)[a-z\s]+(?=\s*)"
    chapterPattern = r"\d+(?=\s*\:)"
    versePattern = r"(?<=\:\s*)[\d]+"
    flags=re.I|re.M

    raw_data = re.search(allBiblePattern, str(text), flags=flags)
    if raw_data is not None:
        raw_data = raw_data.group().strip()
        book = re.search(bookPattern, raw_data, flags=flags)
        chapters = re.findall(chapterPattern, raw_data, flags=flags)
        verses = re.findall(versePattern, raw_data, flags=flags)
        if book is not None and verses is not None and chapters is not None:
            return {
                'book': book.group().lower().strip(),
                'chapters': chapters,
                'verses': verses,
            }
        
    return None

def createBs(bibleList):
    bsFile = open('test.prg', "w+")
    bsFile.write("[Programs]\n")
    for i, data in enumerate(bibleList):
        bookIndex = books.get(data['book'])
        if len(data['chapters']) > 1:
            for chapterIndex in range(len(data['chapters'])):
                bsFile.write('Bm{index}=TB:{bookIndex}:{bookChapter}:{bookVerse}\n'.format(
                index=i,bookIndex=bookIndex,bookChapter=data['chapters'][chapterIndex],bookVerse=data['verses'][chapterIndex]))
                i += 1           
        else:
            bsFile.write('Bm{index}=TB:{bookIndex}:{bookChapter}:{bookVerse}\n'.format(
                index=i,bookIndex=bookIndex,bookChapter=data['chapters'][0],bookVerse=data['verses'][0]))
    bsFile.close()



# App entry point
def main():
    if len(sys.argv) < 2:
        return print("Please provide .pptx file\nExample: python3 create_bs.py ./test.pptx")
    
    ppt = Presentation(sys.argv[1])
    bibleList = []

    # Looping through each slide in ppt and scanning all bible
    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            data = extractBible(repr(shape.text)) # using repr to passing raw text
            if data is not None:
                bibleList.append(data)
                
    # Creating bibleshow from list of bible   
    createBs(bibleList)

main()
